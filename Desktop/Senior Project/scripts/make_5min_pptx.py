#!/usr/bin/env python3
"""
Generate a 5-minute PowerPoint presentation for the Dark Fleets project.
Requires: pip install python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# ==================== CONFIG ====================
prs = Presentation()
prs.slide_width = Inches(13.33)   # 16:9
prs.slide_height = Inches(7.5)

# Dark theme colors
BG_COLOR = RGBColor(15, 23, 42)      # #0f172a
TEXT_COLOR = RGBColor(226, 232, 240)  # #e2e8f0
ACCENT_COLOR = RGBColor(59, 130, 246)  # #3b82f6


def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background(slide)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(54)
    title_shape.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR


def add_content_slide(title, bullets, image_path=None, image_left=Inches(7), image_top=Inches(1.5), image_width=Inches(5.5)):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide)
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR

    tf = slide.placeholders[1].text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = bullets[0]
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT_COLOR
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_COLOR

    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, image_left, image_top, width=image_width)
        except Exception:
            pass  # image missing → skip gracefully


# ==================== 5-MIN SLIDES ====================

add_title_slide(
    "Graph-Based Detection\nof Cooperative Dark Fleets",
    "5-Minute Presentation\n[Your Name] • February 2026",
)

add_content_slide(
    "Motivation & Goal",
    [
        "IUU fishing fleets coordinate at sea",
        "• Transshipment",
        "• Shared fishing grounds",
        '• Rendezvous with "dark" vessels (AIS off)',
        "",
        "Goal: Surface candidate pairs WITHOUT labels",
        "→ Unsupervised link prediction + validation",
    ],
)

add_content_slide(
    "Approach (High-Level)",
    [
        "1. Daily AIS data 2012–2019 (2.38 billion records)",
        "2. Build temporal graph (edge = <10 km same day)",
        "3. Train Temporal GCN (TGCN) + temporal node features",
        "4. Score missing links → multi-threshold validation",
        "",
        "25 km ±1d | 50 km ±3d | 100 km ±7d",
    ],
)

add_content_slide(
    "Results",
    [
        "Full-coverage model: ROC AUC ≈ 0.78",
        "Capped reproducible: ROC AUC ≈ 0.95",
        "5-fold CV: 0.927 ± 0.067",
        "",
        "Strongest pair: 412000690 ↔ 412325200",
        "→ 102 days within 25 km",
    ],
    image_path="artifacts/plots/case_study_pairs/pair_412000690_412325200.png",
)

add_content_slide(
    "Conclusion & Future Work",
    [
        "Reproducible pipeline that surfaces credible candidates",
        "Ready for enforcement & research",
        "",
        "Next steps:",
        "• Finer resolution",
        "• Vessel metadata lookup",
        "• Supervised learning (if labels available)",
        "",
        "Thank you! Questions?",
    ],
    image_path=None,
)

# Save
out_path = "docs/5_Minute_Dark_Fleets_Presentation.pptx"
os.makedirs(os.path.dirname(out_path), exist_ok=True)
prs.save(out_path)
print(f"✅ 5-minute PPTX created: {out_path}")
