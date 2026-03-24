#!/usr/bin/env python3
"""
Generate a 20-minute PowerPoint presentation for the Dark Fleets project.
Requires: pip install python-pptx

To ensure all 4 case study plots exist (including pair 978925333/415131223):
  python3 scripts/compute_pair_overlap_from_daily.py --pairs artifacts/case_study_pairs.csv \\
    --daily-root "data/MMSI daily vessels " --top-k 4 --distance-km 25 --day-window 1 \\
    --max-files-per-year 100 --out-dir artifacts/plots/case_study_pairs
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# Same config as 5-min script
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BG_COLOR = RGBColor(15, 23, 42)
TEXT_COLOR = RGBColor(226, 232, 240)
ACCENT_COLOR = RGBColor(59, 130, 246)


def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background(slide)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(54)
    title_shape.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR


def add_content_slide(title, bullets, image_path=None, image_left=Inches(7), image_top=Inches(1.5), image_width=Inches(5.5)):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide)
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR

    tf = slide.placeholders[1].text_frame
    tf.clear()
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24) if "•" in bullet else Pt(28)
        p.font.color.rgb = TEXT_COLOR

    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, image_left, image_top, width=image_width)
        except Exception:
            pass


# ==================== 20-MIN SLIDES ====================

add_title_slide(
    "Graph-Based Detection\nof Cooperative Dark Fleets",
    "20-Minute Presentation\n[Your Name] • February 2026",
)

add_content_slide(
    "Agenda",
    [
        "1. Introduction & Motivation (3 min)",
        "2. Data Pipeline (3 min)",
        "3. Methods (5 min)",
        "4. Results & Validation (5 min)",
        "5. Case Studies (3 min)",
        "6. Limitations & Conclusion (1 min)",
    ],
)

add_content_slide(
    "1. Motivation",
    [
        "IUU fishing costs billions",
        "Key pattern: vessel coordination at sea",
        "• Transshipment",
        "• Shared fishing grounds",
        "• Rendezvous with dark vessels",
        "",
        "ZERO labeled pairs → Unsupervised approach",
    ],
)

add_content_slide(
    "2. Data Pipeline",
    [
        "2.38 billion AIS records (2012–2019)",
        "SST coverage: 97.5%",
        "",
        "Pipeline: Chunked merge → Year split → SST lookup → QC → Temporal graph",
        "Edge = vessels < 10 km same day (876 daily snapshots)",
    ],
)

add_content_slide(
    "3. Methods – Graph & TGCN",
    [
        "Nodes = MMSI (hundreds of thousands)",
        "Temporal node features: interactions, partners, last seen, mean gap",
        "",
        "TGCN = Graph Conv + GRU",
        "Score = dot-product of embeddings",
    ],
)

add_content_slide(
    "3. Methods – Evaluation",
    [
        "Strict time split (past → future)",
        "5-fold rolling-window CV (50%–90% train)",
        "",
        "Metrics: ROC AUC + Average Precision",
        "3 random seeds",
    ],
)

add_content_slide(
    "4. Results – Performance",
    [
        "Full-coverage (1500 buckets): ≈ 0.78",
        "Capped graph: ≈ 0.95",
        "5-fold CV: 0.927 ± 0.067",
        "",
        "Temporal features + TGCN clearly outperform static baselines",
    ],
)

add_content_slide(
    "4. Results – Validation",
    [
        "Top-1000 candidates ranked",
        "",
        "Multi-threshold check:",
        "• 25 km ±1 day → 8 pairs",
        "• 50 km ±3 days → 11 pairs",
        "• 100 km ±7 days → 14 pairs",
        "",
        'Filters "same region" from real close encounters',
    ],
)

add_content_slide(
    "Understanding the Numbers",
    [
        "Proximity days (e.g. 102) = days we found the pair within 25 km in GPS data",
        "Candidate score (e.g. 15.92) = model's confidence; higher = stronger",
        "",
        "These appear on the next slides.",
    ],
)

# Case Study Slides with the 4 graphs
add_content_slide(
    "Case Study 1 – Strongest Candidate",
    [
        "412000690 ↔ 412325200",
        "Score: 15.92",
        "102 days within 25 km (±1 day)",
        "410 days within 50 km (±3 days)",
        "1,568 days within 100 km (±7 days)",
        "Mean distance ~135–139 km",
        "Multi-year repeated close encounters",
        "→ Highest-confidence cooperative candidate",
    ],
    image_path="artifacts/plots/case_study_pairs/pair_412000690_412325200.png",
)

add_content_slide(
    "Case Study 2 – Moderate Evidence",
    [
        "412061791 ↔ 412508302",
        "Score: 15.81",
        "8 days within 25 km (±1 day)",
        "88 days within 50 km (±3 days)",
        "378 days within 100 km (±7 days)",
        "Intermittent close proximity",
        "Possible transshipment or shared grounds",
        "→ Solid secondary candidate",
    ],
    image_path="artifacts/plots/case_study_pairs/pair_412061791_412508302.png",
)

add_content_slide(
    "Case Study 3 – Occasional Coordination",
    [
        "412425192 ↔ 998508450",
        "Score: 14.86",
        "9 days within 25 km (±1 day)",
        "20 days within 50 km (±3 days)",
        "53 days within 100 km (±7 days)",
        "Fewer total encounters",
        "Possible occasional meetings or shared region",
        "→ Useful for geographic diversity",
    ],
    image_path="artifacts/plots/case_study_pairs/pair_412425192_998508450.png",
)

add_content_slide(
    "Case Study 4 – Regional Overlap (Contrast)",
    [
        "978925333 ↔ 415131223",
        "Score: 14.80",
        "2 days within 25 km (±1 day)",
        "12 days within 50 km (±3 days)",
        "864 days within 100 km (±7 days)",
        "Very few close meetings",
        "Likely same fishing grounds / corridor",
        "→ Shows value of multi-threshold validation",
    ],
    image_path="artifacts/plots/case_study_pairs/pair_978925333_415131223.png",
)

add_content_slide(
    "6. Limitations",
    [
        "Daily grid-cell resolution (fine events blurred)",
        "No ground-truth labels (proximity = proxy)",
        "Results depend on thresholds & time period",
    ],
)

add_content_slide(
    "Conclusion & Future Work",
    [
        "Working, reproducible pipeline",
        "Strongest candidate ready for investigation",
        "",
        "Future:",
        "• Sub-daily resolution",
        "• Metadata (flag, gear, ownership)",
        "• Supervised learning when labels available",
        "",
        "Thank you! Questions?",
    ],
)

# Save
out_path = "docs/20_Minute_Dark_Fleets_Presentation.pptx"
os.makedirs(os.path.dirname(out_path), exist_ok=True)
prs.save(out_path)
print(f"✅ 20-minute PPTX created: {out_path}")
